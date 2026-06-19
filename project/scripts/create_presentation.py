from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

from src.currency import USD_TO_KRW_RATE
from src.data_loader import load_used_cars
from src.model import train_price_model


OUTPUT_PATH = PROJECT_ROOT / "중고차_데이터_분석_발표자료.pptx"
SCRIPT_PATH = PROJECT_ROOT / "docs" / "presentation_script.md"

WIDE_WIDTH = Inches(13.333)
WIDE_HEIGHT = Inches(7.5)

INK = RGBColor(24, 32, 42)
MUTED = RGBColor(91, 103, 112)
PAPER = RGBColor(248, 247, 243)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(214, 210, 200)
ACCENT = RGBColor(25, 118, 210)
ACCENT_2 = RGBColor(232, 135, 65)
GREEN = RGBColor(47, 140, 92)
RED = RGBColor(196, 70, 70)
DEEP = RGBColor(16, 35, 55)


def money_krw(value: float) -> str:
    if value >= 100_000_000:
        return f"{value / 100_000_000:.1f}억원"
    return f"{value / 10_000:.0f}만원"


def add_bg(slide, color=PAPER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, WIDE_WIDTH, WIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=18, color=INK, bullet=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        if bullet:
            p.level = 0
            p.text = f"• {line}"
    return box


def add_kicker(slide, text, page):
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.38), Inches(0.12), Inches(0.12))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT_2
    marker.line.fill.background()
    add_text(slide, text.upper(), Inches(0.82), Inches(0.31), Inches(3.2), Inches(0.28), size=10, bold=True, color=MUTED)
    add_text(slide, f"{page:02d}", Inches(12.05), Inches(0.3), Inches(0.6), Inches(0.25), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, kicker, title, subtitle, page):
    add_kicker(slide, kicker, page)
    add_text(slide, title, Inches(0.62), Inches(0.78), Inches(8.9), Inches(0.85), size=34, bold=True, color=INK)
    if subtitle:
        add_text(slide, subtitle, Inches(0.65), Inches(1.52), Inches(9.7), Inches(0.45), size=15, color=MUTED)


def add_metric(slide, x, y, w, h, value, label, note=None, fill=PANEL, value_color=INK):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    box.line.width = Pt(1)
    add_text(slide, value, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.35), size=23, bold=True, color=value_color)
    add_text(slide, label, x + Inches(0.18), y + Inches(0.58), w - Inches(0.36), Inches(0.28), size=11, bold=True, color=MUTED)
    if note:
        add_text(slide, note, x + Inches(0.18), y + Inches(0.9), w - Inches(0.36), Inches(0.32), size=10, color=MUTED)
    return box


def add_footer(slide):
    add_text(
        slide,
        "Source: Kaggle - Used Car Price Prediction Dataset by Taeef Najib\nAdditional data: CARFAX sample, EPA vehicle data, NHTSA recall summary",
        Inches(0.64),
        Inches(6.94),
        Inches(12.0),
        Inches(0.36),
        size=7,
        color=MUTED,
    )


def style_chart(chart):
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.name = "Malgun Gothic"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = "Malgun Gothic"
    chart.chart_title.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ACCENT


def add_bar_chart(slide, df, x, y, w, h, title):
    data = CategoryChartData()
    data.categories = list(df.index)
    data.add_series("평균 가격 (백만원)", list(df.round(0).values))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    style_chart(chart)
    chart.value_axis.tick_labels.number_format = '#,##0,,"백만원"'
    return chart


def add_line_chart(slide, series, x, y, w, h, title):
    data = CategoryChartData()
    data.categories = [str(int(v)) for v in series.index]
    data.add_series("평균 가격 (백만원)", list(series.round(0).values))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = False
    chart.value_axis.tick_labels.number_format = '#,##0,,"백만원"'
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.line.color.rgb = ACCENT_2
    chart.series[0].format.line.width = Pt(2.5)
    return chart


def add_flow(slide, steps, x, y, box_w, box_h, gap):
    last = None
    for index, (head, body, color) in enumerate(steps):
        bx = x + index * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, bx, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = LINE
        add_text(slide, head, bx + Inches(0.14), y + Inches(0.14), box_w - Inches(0.28), Inches(0.25), size=13, bold=True, color=INK)
        add_text(slide, body, bx + Inches(0.14), y + Inches(0.48), box_w - Inches(0.28), box_h - Inches(0.58), size=10, color=MUTED)
        if last is not None:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, last, y + box_h / 2, bx, y + box_h / 2)
            conn.line.color.rgb = MUTED
            conn.line.width = Pt(1.5)
        last = bx + box_w


def create_deck():
    df = load_used_cars()
    _, metrics = train_price_model(df)
    original_rows = 4009
    enriched_rows = len(df)
    epa_matches = int(df.get("external_epa_matched", pd.Series(dtype=bool)).sum())
    nhtsa_matches = int(df.get("external_nhtsa_matched", pd.Series(dtype=bool)).sum())
    top_brand_avg = (
        df.groupby("brand")
        .filter(lambda group: len(group) >= 80)
        .groupby("brand")["price"]
        .mean()
        .sort_values(ascending=False)
        .head(6)
        / 1_000_000
    )
    recent_year_avg = (
        df[df["model_year"].between(2015, 2026)]
        .groupby("model_year")["price"]
        .mean()
        .sort_index()
        / 1_000_000
    )
    price_limit_usd = metrics["price_limit"] / USD_TO_KRW_RATE
    mae_usd = metrics["mae"] / USD_TO_KRW_RATE

    prs = Presentation()
    prs.slide_width = WIDE_WIDTH
    prs.slide_height = WIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, DEEP)
    add_text(slide, "BIGDATA FINAL PROJECT", Inches(0.7), Inches(0.5), Inches(4.2), Inches(0.3), size=11, bold=True, color=RGBColor(188, 210, 230))
    add_text(slide, "이 중고차,\n지금 가격이 괜찮은 걸까?", Inches(0.68), Inches(1.2), Inches(8.8), Inches(1.85), size=42, bold=True, color=RGBColor(255, 255, 255))
    add_text(slide, "매물 데이터를 분석하고, 조건별 참고 가격을 확인하는 Streamlit 서비스를 만들었습니다.", Inches(0.72), Inches(3.08), Inches(8.8), Inches(0.45), size=17, color=RGBColor(218, 228, 236))
    add_metric(slide, Inches(0.78), Inches(4.35), Inches(2.7), Inches(1.22), f"{enriched_rows:,}", "확장 데이터 행 수", "Kaggle 4,009행 + 보조 샘플", fill=RGBColor(235, 244, 251), value_color=DEEP)
    add_metric(slide, Inches(3.75), Inches(4.35), Inches(2.7), Inches(1.22), "핵심 4단계", "발표 흐름", "EDA→시각화→예측→추천", fill=RGBColor(255, 241, 229), value_color=DEEP)
    add_metric(slide, Inches(6.72), Inches(4.35), Inches(2.7), Inches(1.22), "R² 0.677", "예측 모델", "RandomForestRegressor", fill=RGBColor(233, 246, 239), value_color=DEEP)
    add_text(slide, "발표자: 프로젝트 제출자", Inches(0.76), Inches(6.65), Inches(3.0), Inches(0.28), size=11, color=RGBColor(218, 228, 236))

    # 2. Problem
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "문제", "중고차 가격은 숫자 하나만 보고 판단하기 어렵습니다.", "같은 브랜드라도 연식, 주행거리, 사고 이력에 따라 가격 차이가 커짐", 2)
    checks = [
        ("비교 기준 부족", "매물 가격이 싼지 비싼지 바로 판단하기 어려움"),
        ("조건 조합 복잡", "브랜드, 모델, 연식, 주행거리가 함께 작용"),
        ("평균의 함정", "초고가 매물이 섞이면 일반 가격대가 흐려짐"),
        ("앱으로 해결", "조건 입력 → 예측 가격 → 유사 매물 비교"),
    ]
    for i, (head, body) in enumerate(checks):
        y = Inches(2.25 + i * 0.92)
        add_metric(slide, Inches(0.8), y, Inches(2.15), Inches(0.68), "✓", head, None, fill=RGBColor(234, 244, 238), value_color=GREEN)
        add_text(slide, body, Inches(3.15), y + Inches(0.16), Inches(7.6), Inches(0.28), size=18, color=INK)
    add_multiline(
        slide,
        ["핵심 질문: 입력한 차량 조건이 시장 데이터 안에서 어느 정도 가격대로 보이는가?", "발표 흐름: 데이터 준비 → EDA 발견 → 예측 모델 → 앱 시연"],
        Inches(0.85),
        Inches(6.05),
        Inches(11.1),
        Inches(0.75),
        size=14,
        color=MUTED,
    )

    # 3. Data sources
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "데이터", "Kaggle 중고차 매물 데이터를 기본으로 사용하고, 최신 연식 경향 참고용 샘플과 외부 정보를 보조로 붙였습니다.", "가격 예측은 매물 조건 중심, 추천과 판독은 연비·리콜 정보까지 참고", 3)
    add_metric(slide, Inches(0.8), Inches(2.1), Inches(2.55), Inches(1.25), f"{original_rows:,}", "Kaggle 원본", "Used Car Price Prediction")
    add_metric(slide, Inches(3.65), Inches(2.1), Inches(2.55), Inches(1.25), f"{enriched_rows:,}", "확장 데이터", "앱에서 사용")
    add_metric(slide, Inches(6.5), Inches(2.1), Inches(2.55), Inches(1.25), f"{epa_matches:,}", "EPA 매칭", "연비/연료 비용")
    add_metric(slide, Inches(9.35), Inches(2.1), Inches(2.55), Inches(1.25), f"{nhtsa_matches:,}", "NHTSA 매칭", "리콜 위험 참고")
    add_flow(
        slide,
        [
            ("Kaggle 매물", "브랜드, 모델, 연식, 주행거리, 가격", RGBColor(245, 247, 250)),
            ("CARFAX 샘플", "최신 연식 경향 참고용 보조 데이터", RGBColor(235, 244, 251)),
            ("EPA/NHTSA", "브랜드·모델·연식 기준 가능한 항목 매칭", RGBColor(255, 241, 229)),
            ("앱 기능", "가격 예측, 추천, 구입 판독에 사용", RGBColor(233, 246, 239)),
        ],
        Inches(0.8),
        Inches(4.55),
        Inches(2.65),
        Inches(1.05),
        Inches(0.32),
    )
    add_footer(slide)

    # 4. Analysis frame
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "분석 기준", "가격을 세 가지 관점으로 나눠서 봤습니다.", "가격 자체, 차량 조건, 외부 위험 정보를 함께 비교", 4)
    source_items = [
        ("가격", "매물 가격과 분포", "평균보다 일반 가격대와 이상치를 구분"),
        ("차량 조건", "브랜드·모델·연식·주행거리", "가격을 흔드는 핵심 입력값"),
        ("상태 정보", "사고 이력·클린 타이틀", "같은 조건 안에서 가격 차이를 설명"),
        ("외부 정보", "EPA 연비·NHTSA 리콜", "추천과 구입 판독에서 참고"),
    ]
    for i, (name, source, role) in enumerate(source_items):
        x = Inches(0.8 + (i % 2) * 5.9)
        y = Inches(2.15 + (i // 2) * 1.5)
        add_metric(slide, x, y, Inches(5.25), Inches(1.05), name, source, role, fill=PANEL, value_color=ACCENT if i < 2 else GREEN)
    add_text(slide, "발표에서 전제로 둘 부분", Inches(0.9), Inches(5.55), Inches(2.8), Inches(0.28), size=17, bold=True)
    add_multiline(
        slide,
        [
            "실시간 시세 서비스가 아니라 샘플 데이터 기반 참고 서비스",
            "모델 결과는 감정가가 아니라 조건별 기준선으로 해석",
            "매칭되지 않은 외부 정보는 Unknown 또는 결측값으로 처리",
            "한국 기준 보정은 발표 핵심이 아니라 보조 비교 기능",
        ],
        Inches(0.9),
        Inches(5.9),
        Inches(11.4),
        Inches(1.15),
        size=13,
        color=INK,
        bullet=True,
    )

    # 5. Preprocessing
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "분석 준비", "문자로 들어온 매물 정보를 모델이 읽을 수 있게 정리했습니다.", "가격과 주행거리를 숫자로 바꾸고, 사고 여부 같은 판단 변수를 추가", 5)
    add_flow(
        slide,
        [
            ("문자열 정제", "$, 쉼표, mi. 제거 후 숫자형 변환", RGBColor(245, 247, 250)),
            ("통화 통일", "price_usd 보존, 고정 환율로 원화 변환", RGBColor(235, 244, 251)),
            ("결측 처리", "연료, 사고, 타이틀, 변속기 Unknown 처리", RGBColor(255, 241, 229)),
            ("변수 추가", "차량 나이, 주행거리당 가격, 사고 여부", RGBColor(233, 246, 239)),
        ],
        Inches(0.85),
        Inches(2.25),
        Inches(2.7),
        Inches(1.2),
        Inches(0.28),
    )
    add_text(slide, "모델에 넣기 전 정리", Inches(0.95), Inches(4.55), Inches(2.4), Inches(0.3), size=18, bold=True)
    add_multiline(
        slide,
        [
            "연식과 주행거리는 숫자형 변수로 사용",
            "브랜드, 모델, 연료, 변속기 등은 원-핫 인코딩",
            "가격 분포가 한쪽으로 치우쳐 로그 변환 적용",
            "환율 기준: 1 USD = 1,532.52원 (2026-06-04)",
            "상위 1% 초고가 매물 93건은 학습에서 제외하고 따로 표시",
        ],
        Inches(0.95),
        Inches(5.0),
        Inches(10.9),
        Inches(1.48),
        size=14,
        color=INK,
        bullet=True,
    )

    # 6. EDA
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "EDA 발견", "가격은 브랜드와 연식에 따라 확실히 달라졌습니다.", "초고가 매물이 평균을 흔들어서 일반 가격대는 따로 봐야 했음", 6)
    add_bar_chart(slide, top_brand_avg, Inches(0.75), Inches(2.05), Inches(5.55), Inches(4.3), "브랜드별 평균 가격 (백만원)")
    add_line_chart(slide, recent_year_avg, Inches(6.75), Inches(2.05), Inches(5.55), Inches(4.3), "연식별 평균 가격 추세 (백만원)")
    add_multiline(
        slide,
        ["상위 1% 초고가 매물은 따로 분리해 일반 매물과 구분", "최신 연식일수록 가격은 높고, 주행거리가 늘수록 가격은 낮아지는 경향"],
        Inches(0.8),
        Inches(6.55),
        Inches(11.6),
        Inches(0.55),
        size=12,
        color=MUTED,
    )

    # 7. Model
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "예측 모델", "모델의 역할은 정답 맞히기가 아니라 비교 기준을 주는 것입니다.", "예측 가격과 유사 매물 평균을 함께 보여줘서 결과를 해석하게 구성", 7)
    add_metric(slide, Inches(0.85), Inches(1.95), Inches(3.05), Inches(1.18), money_krw(metrics["mae"]), "평균 절대 오차", f"약 ${mae_usd:,.0f}")
    add_metric(slide, Inches(4.15), Inches(1.95), Inches(3.05), Inches(1.18), f"{metrics['r2']:.3f}", "R² 설명력", "가격 변동의 일부 설명")
    add_metric(slide, Inches(7.45), Inches(1.95), Inches(3.05), Inches(1.18), f"{metrics['excluded_rows']:,}", "학습 제외 매물", f"상위 1% 기준 ${price_limit_usd:,.0f}+")
    add_text(slide, "RandomForest를 쓴 이유", Inches(0.95), Inches(3.55), Inches(3.2), Inches(0.28), size=17, bold=True)
    add_multiline(
        slide,
        [
            "브랜드와 모델처럼 범주형 변수가 많음",
            "연식, 주행거리, 사고 이력 조합에 따라 가격이 달라짐",
            "조건 조합별 가격 패턴을 비교적 유연하게 반영",
        ],
        Inches(0.95),
        Inches(3.95),
        Inches(5.2),
        Inches(1.1),
        size=13,
        color=INK,
        bullet=True,
    )
    add_text(slide, "결과 해석", Inches(6.75), Inches(3.55), Inches(2.2), Inches(0.28), size=17, bold=True)
    add_multiline(
        slide,
        [
            "MAE 1,498만원은 작지 않은 오차",
            "R² 0.677로 가격 변동 일부를 설명",
            "최종 결과는 참고 가격대로 보는 것이 적절",
        ],
        Inches(6.75),
        Inches(3.95),
        Inches(5.2),
        Inches(1.1),
        size=13,
        color=INK,
        bullet=True,
    )
    add_text(slide, "발표 포인트: 예측값 하나만 보여주지 않고, 비슷한 차량의 평균/중앙값과 같이 비교", Inches(0.95), Inches(6.15), Inches(10.2), Inches(0.28), size=16, bold=True, color=ACCENT)
    add_footer(slide)

    # 8. Service
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "앱 시연", "분석 결과를 사용자가 직접 눌러보고 비교할 수 있게 했습니다.", "발표에서는 네 화면만 빠르게 보여주면 흐름이 잡힘", 8)
    cols = [
        ("EDA / 시각화", "가격 분포와 조건별 차이를 먼저 확인", ACCENT),
        ("모델 / 예측", "차량 조건을 넣고 예상 가격 확인", GREEN),
        ("추천", "예산 이하에서 최신 연식·낮은 주행거리 우선", ACCENT_2),
        ("구입 판독", "예측가 대비 가격, 사고, 연비, 리콜을 함께 비교", RED),
    ]
    for i, (head, body, color) in enumerate(cols):
        x = Inches(0.85 + (i % 2) * 5.9)
        y = Inches(2.15 + (i // 2) * 1.85)
        add_metric(slide, x, y, Inches(5.25), Inches(1.35), head, body, None, fill=PANEL, value_color=color)
    add_text(slide, "데모 순서: EDA 분포 → 브랜드/연식 필터 → 조건 입력 예측 → 예산 이하 추천", Inches(0.9), Inches(6.3), Inches(11.1), Inches(0.34), size=16, bold=True, color=INK)
    add_text(slide, "보조 기능: 한국 기준 보정은 기준가액 CSV가 있을 때만 쓰는 참고 비교", Inches(0.9), Inches(6.68), Inches(11.4), Inches(0.26), size=12, color=MUTED)

    # 9. Close
    slide = prs.slides.add_slide(blank)
    add_bg(slide, DEEP)
    add_text(slide, "결론", Inches(0.75), Inches(0.55), Inches(1.8), Inches(0.32), size=13, bold=True, color=RGBColor(188, 210, 230))
    add_text(slide, "중고차 가격을 그냥 보는 것에서\n조건별로 비교하고 판단하는 서비스로 만들었습니다.", Inches(0.75), Inches(1.25), Inches(10.9), Inches(1.3), size=34, bold=True, color=RGBColor(255, 255, 255))
    add_multiline(
        slide,
        [
            "구현한 것: EDA, 시각화, 가격 예측, 추천, 구입 판독",
            "확인한 것: 브랜드, 연식, 주행거리, 사고 이력에 따른 가격 차이",
            "한계: 옵션·지역·판매자 정보와 실시간성이 부족함",
            "주의: 해외 데이터 기반이라 한국 원화 가격 해석에는 시장 차이가 있음",
            "개선 방향: 실제 매물, 옵션 데이터, 비교 모델을 추가해 신뢰도 높이기",
        ],
        Inches(0.85),
        Inches(3.2),
        Inches(10.9),
        Inches(1.8),
        size=18,
        color=RGBColor(230, 238, 245),
        bullet=True,
    )
    add_text(slide, "감사합니다", Inches(0.84), Inches(6.35), Inches(3.2), Inches(0.4), size=20, bold=True, color=RGBColor(255, 255, 255))

    prs.save(OUTPUT_PATH)
    write_script()
    return OUTPUT_PATH


def write_script():
    SCRIPT_PATH.write_text(
        """# 중고차 데이터 분석 및 가격 예측 서비스 발표 대본

## 1. 표지
안녕하세요. 저는 중고차 매물 데이터를 활용해 가격 분포를 분석하고, 사용자가 차량 조건을 입력하면 예상 가격을 확인할 수 있는 Streamlit 서비스를 만들었습니다.

## 2. 과제 대응
학생안내 가이드의 핵심은 EDA, 시각화, 입력 기반 예측/서비스까지 동작하는 앱입니다. 제 프로젝트는 이 기본 요구사항에 가격 예측 게임, 차량 추천, 구입 판독 기능까지 확장했습니다. 한국 기준 보정은 미국 매물 중앙값과 한국 기준가액을 비교하는 선택적 보조 기능으로 두고, 발표에서는 EDA, 시각화, 예측, 추천 흐름에 집중하겠습니다.

## 3. 데이터
기본 데이터는 Kaggle의 Used Car Price Prediction Dataset by Taeef Najib이고, 원본 기준 4,009행입니다. 최신 연식 경향을 참고하기 위해 CARFAX 샘플을 추가로 붙였습니다. 다만 실시간 시세 데이터는 아니기 때문에 보조적인 확장 데이터로 해석했고, 현재 앱에서는 9,273행의 확장 데이터를 우선 사용합니다.

## 4. 출처·라이선스
본 프로젝트의 기본 중고차 매물 데이터는 Kaggle의 Used Car Price Prediction Dataset by Taeef Najib을 사용했습니다. 이후 분석과 서비스 기능 확장을 위해 CARFAX 샘플, EPA 연비 데이터, NHTSA 리콜 요약 정보를 일부 보조 데이터로 매칭했습니다. EPA와 NHTSA 데이터는 브랜드, 모델, 연식을 기준으로 가능한 항목만 매칭했고, 매칭되지 않은 차량은 Unknown 또는 결측값으로 처리했습니다.

## 5. 전처리
가격과 주행거리는 문자열이라 기호와 단위를 제거해 숫자로 바꿨습니다. 원본 달러 가격은 price_usd로 보존하고, 앱 기준 가격은 고정 환율 1달러 1,532.52원으로 원화 변환했습니다. 원화 가격은 실제 시세라기보다 발표와 비교를 쉽게 하기 위한 변환값입니다.

## 6. EDA 발견
가격은 일부 초고가 차량 때문에 분포가 크게 왜곡됩니다. 그래서 상위 1% 매물은 비싼차 모음집으로 분리했습니다. Porsche, Land Rover, Cadillac 같은 고가 브랜드가 평균 가격 상위권에 있었고, 연식은 대체로 최신일수록 가격이 높았습니다. 일부 연도는 샘플 구성이나 고가 매물 영향으로 평균이 흔들릴 수 있습니다.

## 7. 모델 결과
모델은 RandomForestRegressor를 사용했습니다. 브랜드와 모델 같은 범주형 변수, 연식과 주행거리 같은 수치형 변수가 섞여 있고 가격 관계가 비선형적이라 이 모델을 선택했습니다. 현재 확장 데이터 기준 평균 절대 오차는 약 1,498만원, R²는 약 0.677입니다. MAE가 작지 않기 때문에 정밀한 감정가 산정보다는 조건별 참고 가격대를 제공하는 모델로 해석했습니다. 향후에는 Linear Regression, XGBoost 등과 비교해 모델 선택 근거를 강화할 필요가 있습니다.

## 8. 서비스 구성
앱은 EDA, 시각화, 모델/서비스, 가격 예측 게임, 추천, 구입 판독으로 구성했습니다. 발표 시연은 EDA 가격 분포, 브랜드/연식별 시각화, 조건 입력 예측, 예산 이하 차량 추천 네 가지를 중심으로 보여드리겠습니다. 게임과 한국 기준 보정은 보조 기능으로 구현했다는 정도만 언급하겠습니다.

## 9. 결론
이 프로젝트는 단순 그래프 출력이 아니라 분석 결과를 사용자가 직접 눌러보고 판단할 수 있는 서비스로 확장한 점이 핵심입니다. 한계는 트림과 옵션, 사고 수리 범위, 지역, 판매자 유형 정보가 부족하다는 점입니다. 또한 Kaggle 기반 샘플 데이터라 실시간 시세와 차이가 있을 수 있고, 해외 데이터 기반이라 한국 원화 가격으로 해석할 때 시장 차이가 있습니다.
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    output = create_deck()
    print(output)
